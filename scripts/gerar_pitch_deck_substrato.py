#!/usr/bin/env python
"""
Gera uma apresentação comercial (PPTX) do SUBSTRATO para potenciais compradores.

Uso:
  python scripts/gerar_pitch_deck_substrato.py

Saída:
  artefatos/SUBSTRATO_PitchDeck.pptx
"""

from __future__ import annotations

from datetime import UTC, datetime
import itertools
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]


def rgb(hex_str: str) -> RGBColor:
    h = hex_str.strip().lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


COLORS = {
    "bg": rgb("#EEF0F2"),
    "white": rgb("#FFFFFF"),
    "text": rgb("#1F1F1F"),
    "muted": rgb("#4B5563"),
    "border": rgb("#D6D9DE"),
    "red": rgb("#A82323"),
    "red_dark": rgb("#7A1B1B"),
    "blue": rgb("#0284C7"),
}


def set_bg(slide, prs: Presentation, color: RGBColor):
    # Background is a full-bleed rectangle created first (so other elements stay above).
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


def add_header(slide, prs: Presentation, title: str):
    h = Inches(0.62)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["red_dark"]
    bar.line.fill.background()

    tf = bar.text_frame
    tf.clear()
    tf.margin_left = Inches(0.45)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.LEFT


def add_footer(slide, prs: Presentation, text: str):
    box = slide.shapes.add_textbox(
        Inches(0.45), prs.slide_height - Inches(0.45), prs.slide_width - Inches(0.9), Inches(0.3)
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = "Calibri"
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["muted"]


def add_bullets(slide, left, top, width, height, title: str, bullets: list[str], title_size=18, bullet_size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)

    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.name = "Calibri"
    p0.font.size = Pt(title_size)
    p0.font.bold = True
    p0.font.color.rgb = COLORS["text"]

    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(bullet_size)
        p.font.color.rgb = COLORS["text"]
        p.space_before = Pt(4)
        p.space_after = Pt(0)

    return box


def add_card(slide, x, y, w, h, title: str, body: str, accent: RGBColor | None = None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["white"]
    card.line.color.rgb = COLORS["border"]
    card.line.width = Pt(1)

    if accent:
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.10), h)
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent
        stripe.line.fill.background()

    tx = slide.shapes.add_textbox(x + Inches(0.20), y + Inches(0.18), w - Inches(0.28), h - Inches(0.25))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = "Calibri"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = COLORS["red_dark"]

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Calibri"
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLORS["muted"]
    p2.space_before = Pt(6)

    return card


def add_big_statement(slide, x, y, w, h, text: str):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = COLORS["text"]
    return box


def add_kpi(slide, x, y, w, h, title: str, body: str, color: RGBColor):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS["white"]
    box.line.color.rgb = COLORS["border"]
    box.line.width = Pt(1)

    tag = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.14))
    tag.fill.solid()
    tag.fill.fore_color.rgb = color
    tag.line.fill.background()

    tx = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.20), w - Inches(0.44), h - Inches(0.26))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = "Calibri"
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = COLORS["text"]

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Calibri"
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLORS["muted"]
    p2.space_before = Pt(8)

    return box


def build_deck() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    img_invoice = ROOT / "artefatos/pdf_preview/fatura_exemplo_preview.png"

    out_dir = ROOT / "artefatos"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "SUBSTRATO_PitchDeck.pptx"

    today = datetime.now(tz=UTC).date().strftime("%Y-%m-%d")

    # Slide 1: Capa
    s1 = prs.slides.add_slide(blank)
    set_bg(s1, prs, COLORS["white"])

    # Left panel with an abstract "network" pattern (avoid embedding third-party branding).
    left_w = Inches(8.0)
    left_panel = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, left_w, prs.slide_height)
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = COLORS["bg"]
    left_panel.line.fill.background()

    nodes = [
        (Inches(0.8), Inches(1.4), Inches(0.30)),
        (Inches(2.0), Inches(0.9), Inches(0.22)),
        (Inches(3.2), Inches(1.6), Inches(0.26)),
        (Inches(4.4), Inches(1.1), Inches(0.20)),
        (Inches(5.7), Inches(1.8), Inches(0.28)),
        (Inches(6.7), Inches(1.2), Inches(0.18)),
        (Inches(7.2), Inches(2.2), Inches(0.24)),
        (Inches(6.0), Inches(2.8), Inches(0.18)),
        (Inches(4.8), Inches(2.4), Inches(0.16)),
        (Inches(3.6), Inches(3.0), Inches(0.22)),
        (Inches(2.2), Inches(2.6), Inches(0.18)),
        (Inches(1.2), Inches(3.1), Inches(0.20)),
    ]

    # Connect nodes
    for (x1, y1, r1), (x2, y2, r2) in itertools.pairwise(nodes):
        conn = s1.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1 + r1, y1 + r1, x2 + r2, y2 + r2)
        conn.line.color.rgb = rgb("#B91C1C")
        conn.line.width = Pt(2)

    # Draw nodes
    for x, y, r in nodes:
        dot = s1.shapes.add_shape(MSO_SHAPE.OVAL, x, y, r * 2, r * 2)
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb("#DC2626")
        dot.line.fill.background()

    panel_x = left_w
    panel_w = prs.slide_width - left_w
    panel = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, panel_x, 0, panel_w, prs.slide_height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["red_dark"]
    panel.line.fill.background()

    # "Logo" container (white capsule)
    logo_box = s1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        panel_x + Inches(0.55),
        Inches(0.55),
        panel_w - Inches(1.10),
        Inches(1.05),
    )
    logo_box.fill.solid()
    logo_box.fill.fore_color.rgb = COLORS["white"]
    logo_box.line.fill.background()

    # python-pptx 0.6.x does not expose MSO_SHAPE.PLUS; use MATH_PLUS instead.
    plus = s1.shapes.add_shape(MSO_SHAPE.MATH_PLUS, panel_x + Inches(0.78), Inches(0.78), Inches(0.34), Inches(0.34))
    plus.fill.solid()
    plus.fill.fore_color.rgb = COLORS["red_dark"]
    plus.line.fill.background()

    logo_text = s1.shapes.add_textbox(panel_x + Inches(1.20), Inches(0.72), panel_w - Inches(1.85), Inches(0.75))
    tf = logo_text.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "SUBSTRATO"
    p.font.name = "Calibri"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["red_dark"]
    p.alignment = PP_ALIGN.LEFT

    title = s1.shapes.add_textbox(panel_x + Inches(0.55), Inches(2.05), panel_w - Inches(1.1), Inches(1.2))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "SUBSTRATO"
    r.font.name = "Calibri"
    r.font.size = Pt(48)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]

    sub = s1.shapes.add_textbox(panel_x + Inches(0.55), Inches(3.10), panel_w - Inches(1.1), Inches(1.1))
    tf = sub.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Infraestrutura unificada de saúde"
    r.font.name = "Calibri"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = rgb("#FECACA")  # light red

    small = s1.shapes.add_textbox(panel_x + Inches(0.55), Inches(4.05), panel_w - Inches(1.1), Inches(1.0))
    tf = small.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Da recepção ao laboratório, faturação e gestão.\nUma única base. Um único fluxo. Um único controlo."
    p.font.name = "Calibri"
    p.font.size = Pt(16)
    p.font.color.rgb = rgb("#FFE4E6")  # very light

    tagline = s1.shapes.add_textbox(
        panel_x + Inches(0.55), prs.slide_height - Inches(1.05), panel_w - Inches(1.1), Inches(0.7)
    )
    tf = tagline.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"SUBS — Sistema Unificado de Base em Saúde  |  Apresentação comercial  |  {today}"
    p.font.name = "Calibri"
    p.font.size = Pt(11)
    p.font.color.rgb = rgb("#FEE2E2")

    # Slide 2: Dor / Problema
    s2 = prs.slides.add_slide(blank)
    set_bg(s2, prs, COLORS["bg"])
    add_header(s2, prs, "O que está a custar caro hoje")

    add_bullets(
        s2,
        Inches(0.7),
        Inches(1.25),
        Inches(7.0),
        Inches(5.7),
        "Problemas reais do dia-a-dia",
        [
            "Processos fragmentados: recepção, laboratório, enfermagem, farmácia e contabilidade não conversam.",
            "Resultados atrasados e sem rastreabilidade: erros, retrabalho e reclamações.",
            "Faturação com fugas: itens esquecidos, IVA inconsistente e pagamentos difíceis de reconciliar.",
            "Falta de auditoria: quem fez o quê, quando e porquê.",
            "Relatórios demorados: decisões no escuro, sem indicadores confiáveis.",
        ],
        title_size=18,
        bullet_size=16,
    )

    add_kpi(
        s2,
        Inches(8.0),
        Inches(1.25),
        Inches(4.7),
        Inches(1.65),
        "Consequência direta",
        "Perdas silenciosas todos os dias: tempo, dinheiro e reputação.",
        COLORS["red"],
    )
    add_kpi(
        s2,
        Inches(8.0),
        Inches(3.10),
        Inches(4.7),
        Inches(1.65),
        "O risco",
        "Sem auditoria e sem fluxo, o erro vira rotina.",
        COLORS["blue"],
    )
    add_kpi(
        s2,
        Inches(8.0),
        Inches(4.95),
        Inches(4.7),
        Inches(2.0),
        "O que o comprador quer",
        "Controle total, documentos profissionais e operação previsível.",
        COLORS["red_dark"],
    )

    add_footer(s2, prs, "SUBSTRATO • Visão comercial")

    # Slide 3: Solução
    s3 = prs.slides.add_slide(blank)
    set_bg(s3, prs, COLORS["white"])
    add_header(s3, prs, "A solução: SUBSTRATO")

    add_big_statement(
        s3,
        Inches(0.8),
        Inches(1.35),
        Inches(12.0),
        Inches(1.6),
        "O sistema operativo da sua unidade de saúde.",
    )

    box = s3.shapes.add_textbox(Inches(0.8), Inches(2.75), Inches(12.0), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = (
        "Unificamos atendimento, exames, resultados, faturação e controlo em um fluxo único "
        "com permissões por função e rastreabilidade."
    )
    p.font.name = "Calibri"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS["muted"]

    add_card(
        s3,
        Inches(0.8),
        Inches(4.2),
        Inches(3.95),
        Inches(2.55),
        "Velocidade",
        "Menos filas, menos retrabalho.\nFluxos claros do início ao fim.",
        COLORS["red"],
    )
    add_card(
        s3,
        Inches(4.9),
        Inches(4.2),
        Inches(3.95),
        Inches(2.55),
        "Controlo",
        "Auditoria real e documentos padronizados.\nSem furos na faturação.",
        COLORS["blue"],
    )
    add_card(
        s3,
        Inches(9.0),
        Inches(4.2),
        Inches(3.85),
        Inches(2.55),
        "Escala",
        "Multi-clínica (multi-tenant) e pronto para cloud/on-prem.\nCresce sem recomeçar.",
        COLORS["red_dark"],
    )

    add_footer(s3, prs, "SUBSTRATO • Unificação operacional")

    # Slide 4: Módulos
    s4 = prs.slides.add_slide(blank)
    set_bg(s4, prs, COLORS["bg"])
    add_header(s4, prs, "Módulos prontos, integrados e por sector")

    cards = [
        ("Recepção", "Check-in, pacientes, requisições, faturas e recibos.", COLORS["red"]),
        ("Laboratório", "Lançar, gravar e validar resultados com rastreio.", COLORS["blue"]),
        ("Medicina", "Consultas, requisições médicas e prontuário (Cardex).", COLORS["red_dark"]),
        ("Enfermagem/Enfermaria", "Sinais vitais, procedimentos, camas e internamentos.", COLORS["red"]),
        ("Farmácia", "Produtos, lotes, estoque (FEFO) e vendas.", COLORS["blue"]),
        ("Faturamento", "Faturas multi-origem, itens e IVA por item.", COLORS["red_dark"]),
        ("Pagamentos", "Pagamentos, conciliação e recibo automático ao liquidar.", COLORS["red"]),
        ("Contabilidade + RH", "Contas, lançamentos, funcionários e escalas.", COLORS["blue"]),
        ("Monitoramento", "Logs, auditoria e métricas para operação previsível.", COLORS["red_dark"]),
    ]

    x0 = Inches(0.55)
    y0 = Inches(1.25)
    w = Inches(4.0)
    h = Inches(1.65)
    gx = Inches(0.2)
    gy = Inches(0.25)

    for i, (t, b, c) in enumerate(cards):
        row = i // 3
        col = i % 3
        x = x0 + col * (w + gx)
        y = y0 + row * (h + gy)
        add_card(s4, x, y, w, h, t, b, c)

    add_footer(s4, prs, "SUBSTRATO • Módulos integrados")

    # Slide 5: Fluxo
    s5 = prs.slides.add_slide(blank)
    set_bg(s5, prs, COLORS["white"])
    add_header(s5, prs, "Fluxo ponta-a-ponta: sem buracos na operação")

    steps = [
        ("1", "Check-in", "Recepção"),
        ("2", "Requisição", "Exames/Consulta"),
        ("3", "Execução", "Lab/Enfermagem/Farmácia"),
        ("4", "Resultado", "Lançar → Gravar → Validar"),
        ("5", "Fatura", "Itens + IVA por item"),
        ("6", "Pagamento", "Total/Parcial + reconciliação"),
        ("7", "Recibo", "PDF com QR/Barcode + histórico"),
    ]

    step_w = Inches(1.55)
    step_h = Inches(1.15)
    gap = Inches(0.14)
    used = (len(steps) * step_w) + ((len(steps) - 1) * gap)
    left = (prs.slide_width - used) / 2
    top = Inches(2.0)
    badge_size = Inches(0.32)

    for i, (n, t, s) in enumerate(steps):
        x = left + i * (step_w + gap)
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, step_w, step_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["bg"]
        card.line.color.rgb = COLORS["border"]
        card.line.width = Pt(1)

        badge = s5.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.10), top + Inches(0.10), badge_size, badge_size)
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLORS["red_dark"]
        badge.line.fill.background()
        tfb = badge.text_frame
        tfb.clear()
        p = tfb.paragraphs[0]
        p.text = n
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]

        tx = s5.shapes.add_textbox(x + Inches(0.12), top + Inches(0.48), step_w - Inches(0.24), step_h - Inches(0.56))
        tf = tx.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = t
        p1.font.name = "Calibri"
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLORS["text"]
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = s
        p2.font.name = "Calibri"
        p2.font.size = Pt(9)
        p2.font.color.rgb = COLORS["muted"]
        p2.alignment = PP_ALIGN.CENTER

        if i < len(steps) - 1:
            ax = x + step_w
            arrow = s5.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE,
                ax + Inches(0.02),
                top + (step_h / 2) - Inches(0.16),
                gap - Inches(0.04),
                Inches(0.32),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS["red"]
            arrow.line.fill.background()

    note = s5.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(12.0), Inches(2.4))
    tf = note.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Tudo registado e auditável: quem fez, quando fez, e o estado de cada etapa."
    p.font.name = "Calibri"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS["red_dark"]

    add_footer(s5, prs, "SUBSTRATO • Fluxo integrado")

    # Slide 6: Documentos (PDF)
    s6 = prs.slides.add_slide(blank)
    set_bg(s6, prs, COLORS["bg"])
    add_header(s6, prs, "Documentos profissionais (PDF) + rastreio")

    add_bullets(
        s6,
        Inches(0.7),
        Inches(1.35),
        Inches(6.0),
        Inches(5.8),
        "O que o seu cliente vê",
        [
            "Faturas, recibos, requisições e resultados com layout institucional.",
            "QR Code + Código de barras (Code128) para rastreio rápido.",
            "Pronto para impressão e arquivamento: claro, compacto e legível.",
            "Menos erro humano. Mais confiança no atendimento.",
        ],
        title_size=18,
        bullet_size=16,
    )

    if img_invoice.exists():
        s6.shapes.add_picture(str(img_invoice), Inches(7.0), Inches(1.25), width=Inches(5.6), height=Inches(6.0))

    add_footer(s6, prs, "SUBSTRATO • PDFs e rastreabilidade")

    # Slide 7: Segurança
    s7 = prs.slides.add_slide(blank)
    set_bg(s7, prs, COLORS["white"])
    add_header(s7, prs, "Segurança, permissões e auditoria real")

    bullets = [
        "RBAC por sector: Recepção, Laboratório, Enfermagem, Farmácia, Contabilidade, RH, Medicina.",
        "Admin apenas para Administrador (controle total de acesso).",
        "Autenticação JWT (acesso/refresh) e isolamento por inquilino (multi-tenant).",
        "Auditoria de atividades e captura de erros para rastreabilidade operacional.",
        "Backups automatizados + health checks (live/ready) + métricas (Prometheus).",
    ]
    add_bullets(
        s7,
        Inches(0.8),
        Inches(1.35),
        Inches(12.0),
        Inches(5.8),
        "Em saúde, controlo não é luxo. É requisito.",
        bullets,
        title_size=20,
        bullet_size=17,
    )
    add_footer(s7, prs, "SUBSTRATO • Segurança e auditoria")

    # Slide 8: Integrações
    s8 = prs.slides.add_slide(blank)
    set_bg(s8, prs, COLORS["bg"])
    add_header(s8, prs, "Integrações e automação (pronto para crescer)")

    add_bullets(
        s8,
        Inches(0.7),
        Inches(1.35),
        Inches(12.6),
        Inches(5.8),
        "Conectividade que vende",
        [
            "Integração com equipamentos: worklist e inbox de resultados (HTTP JSON, chave por equipamento).",
            "Notificações com templates e logs (e-mail/WhatsApp/SMS) com idempotência por referência.",
            "API documentada (OpenAPI) para parceiros, ERPs e integrações futuras.",
            "Processamento assíncrono com Celery: tarefas pesadas sem travar o atendimento.",
        ],
        title_size=20,
        bullet_size=17,
    )
    add_footer(s8, prs, "SUBSTRATO • Integrações")

    # Slide 9: Analytics
    s9 = prs.slides.add_slide(blank)
    set_bg(s9, prs, COLORS["white"])
    add_header(s9, prs, "Dashboards e relatórios para gestão")

    add_bullets(
        s9,
        Inches(0.7),
        Inches(1.35),
        Inches(6.3),
        Inches(5.8),
        "Decisão com dados",
        [
            "KPIs por sector: atendimento, TAT de resultados, faturação, pagamentos, estoque.",
            "Tabelas filtráveis e indicadores consistentes.",
            "Exportação: PDF | CSV | Word.",
        ],
        title_size=20,
        bullet_size=17,
    )

    # Simple bar chart mock
    chart_x = Inches(7.4)
    chart_y = Inches(2.15)
    chart_w = Inches(5.4)
    chart_h = Inches(4.5)
    frame = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, chart_x, chart_y, chart_w, chart_h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = COLORS["bg"]
    frame.line.color.rgb = COLORS["border"]
    frame.line.width = Pt(1)

    labels = ["Atendimentos", "Resultados", "Faturação", "Pagamentos"]
    values = [0.55, 0.75, 0.62, 0.88]
    bar_w = Inches(0.85)
    base_y = chart_y + chart_h - Inches(0.65)
    for i, (lab, v) in enumerate(zip(labels, values, strict=False)):
        x = chart_x + Inches(0.55) + i * Inches(1.2)
        h = Inches(3.0) * v
        y = base_y - h
        bar = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, bar_w, h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS["red"] if i % 2 == 0 else COLORS["blue"]
        bar.line.fill.background()

        t = s9.shapes.add_textbox(x - Inches(0.1), base_y + Inches(0.05), bar_w + Inches(0.2), Inches(0.5))
        tf = t.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = lab
        p.font.name = "Calibri"
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS["muted"]
        p.alignment = PP_ALIGN.CENTER

    add_footer(s9, prs, "SUBSTRATO • Estatísticas")

    # Slide 10: Implementação
    s10 = prs.slides.add_slide(blank)
    set_bg(s10, prs, COLORS["bg"])
    add_header(s10, prs, "Implementação rápida, operação estável")

    add_bullets(
        s10,
        Inches(0.7),
        Inches(1.35),
        Inches(6.3),
        Inches(5.8),
        "Pronto para produção",
        [
            "On-premise ou cloud.",
            "Docker Compose para arranque rápido; Kubernetes para escala.",
            "Postgres + Redis + Celery: robustez de nível empresarial.",
            "Arquitetura modular (Django/DRF + Next.js): fácil de manter e evoluir.",
        ],
        title_size=20,
        bullet_size=17,
    )

    # Mini architecture diagram
    ax = Inches(7.4)
    ay = Inches(1.7)
    box_w = Inches(5.4)
    arch = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ax, ay, box_w, Inches(5.2))
    arch.fill.solid()
    arch.fill.fore_color.rgb = COLORS["white"]
    arch.line.color.rgb = COLORS["border"]
    arch.line.width = Pt(1)

    def arch_box(x, y, w, h, text, color):
        b = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        b.fill.solid()
        b.fill.fore_color.rgb = color
        b.line.fill.background()
        tf = b.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]
        return b

    fx = ax + Inches(0.55)
    fy = ay + Inches(0.65)
    arch_box(fx, fy, Inches(1.7), Inches(0.75), "Frontend\n(Next.js)", COLORS["blue"])
    arch_box(fx + Inches(2.05), fy, Inches(1.7), Inches(0.75), "API\n(Django/DRF)", COLORS["red_dark"])
    arch_box(fx + Inches(0.2), fy + Inches(1.35), Inches(1.7), Inches(0.75), "Postgres", COLORS["red"])
    arch_box(fx + Inches(2.25), fy + Inches(1.35), Inches(1.7), Inches(0.75), "Redis", COLORS["blue"])
    arch_box(fx + Inches(1.2), fy + Inches(2.25), Inches(2.0), Inches(0.75), "Celery\n(Tarefas)", COLORS["red_dark"])
    arch_box(
        fx + Inches(0.2),
        fy + Inches(3.2),
        Inches(3.75),
        Inches(0.75),
        "Integrações\n(E-mail/WhatsApp/Equipamentos)",
        COLORS["red"],
    )

    add_footer(s10, prs, "SUBSTRATO • Deploy e escala")

    # Slide 11: ROI
    s11 = prs.slides.add_slide(blank)
    set_bg(s11, prs, COLORS["white"])
    add_header(s11, prs, "Resultado prático: mais controlo, mais receita, menos caos")

    add_card(
        s11,
        Inches(0.8),
        Inches(1.55),
        Inches(4.1),
        Inches(5.6),
        "Velocidade",
        "Menos filas.\nMenos retrabalho.\nAtendimento previsível.\nResultados com fluxo claro.",
        COLORS["blue"],
    )
    add_card(
        s11,
        Inches(4.95),
        Inches(1.55),
        Inches(4.1),
        Inches(5.6),
        "Receita",
        "Itens certos na fatura.\nIVA por item.\nPagamentos reconciliados.\nRecibo automático ao liquidar.",
        COLORS["red"],
    )
    add_card(
        s11,
        Inches(9.1),
        Inches(1.55),
        Inches(3.95),
        Inches(5.6),
        "Confiança",
        "Auditoria real.\nDocumentos profissionais.\nRastreio por QR/Barcode.\nGestão por indicadores.",
        COLORS["red_dark"],
    )

    add_footer(s11, prs, "SUBSTRATO • Proposta de valor")

    # Slide 12: CTA
    s12 = prs.slides.add_slide(blank)
    set_bg(s12, prs, COLORS["red_dark"])

    t = s12.shapes.add_textbox(Inches(0.85), Inches(1.1), Inches(11.8), Inches(1.0))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Próximo passo: demonstração e plano de implementação"
    p.font.name = "Calibri"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]

    steps_box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(2.35), Inches(11.8), Inches(3.0))
    steps_box.fill.solid()
    steps_box.fill.fore_color.rgb = rgb("#111827")  # near-black slate
    steps_box.line.fill.background()
    steps_box.fill.fore_color.transparency = 0.10

    tx = s12.shapes.add_textbox(Inches(1.25), Inches(2.65), Inches(11.0), Inches(2.5))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = "1) Demo (60 min) com o seu fluxo real"
    p0.font.name = "Calibri"
    p0.font.size = Pt(22)
    p0.font.bold = True
    p0.font.color.rgb = COLORS["white"]
    for line in [
        "2) Piloto (7-14 dias) com dados reais e validação de processos",
        "3) Go-live + treinamento + suporte contínuo",
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.color.rgb = rgb("#E5E7EB")
        p.space_before = Pt(10)

    contact = s12.shapes.add_textbox(Inches(0.85), prs.slide_height - Inches(1.15), Inches(11.8), Inches(0.7))
    tf = contact.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Contacto: substratosys@gmail.com | WhatsApp: +258 84 777 8476"
    p.font.name = "Calibri"
    p.font.size = Pt(14)
    p.font.color.rgb = rgb("#FEE2E2")

    prs.save(str(out_path))
    return out_path


if __name__ == "__main__":
    path = build_deck()
